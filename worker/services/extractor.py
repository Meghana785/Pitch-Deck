"""
Text extraction service for PitchReady worker.

Supports PDF (via PyMuPDF) and PPTX (via python-pptx).
Raises ValueError for unsupported file types or empty extractions.
"""

from __future__ import annotations

import re
from io import BytesIO


def extract_text(file_bytes: bytes, filename: str) -> str:
    """
    Extract plain text from a PDF or PPTX file.

    Args:
        file_bytes: Raw bytes of the uploaded file.
        filename:   Original filename used to detect file type.

    Returns:
        Cleaned, concatenated text string.

    Raises:
        ValueError: If the file type is unsupported or extracted text is < 100 chars.
    """
    lower = filename.lower()

    if lower.endswith(".pdf"):
        return _extract_pdf(file_bytes)
    elif lower.endswith(".pptx"):
        return _extract_pptx(file_bytes)
    else:
        raise ValueError(f"Unsupported file type: '{filename}'. Only .pdf and .pptx are accepted.")


def _extract_pdf(file_bytes: bytes) -> str:
    """
    Extract text from a PDF using PyMuPDF (fitz).

    Concatenates all pages, normalises excessive whitespace.
    """
    import fitz  # PyMuPDF — imported lazily so the module is optional at import time

    doc = fitz.open(stream=file_bytes, filetype="pdf")
    pages: list[str] = []

    for page in doc:
        pages.append(page.get_text("text"))

    raw = "\n".join(pages)

    # Normalise multiple consecutive blank lines to a single blank line
    text = re.sub(r"\n{3,}", "\n\n", raw).strip()

    if len(text) < 100:
        raise ValueError(
            f"PDF extraction yielded only {len(text)} characters — "
            "file may be image-only or empty."
        )

    return text


def _extract_pptx(file_bytes: bytes) -> str:
    """
    Extract text from a PPTX using python-pptx.

    Prepends '--- Slide N ---' before each slide's text block.
    Concatenates all slides with double newlines.
    """
    from pptx import Presentation  # python-pptx

    prs = Presentation(BytesIO(file_bytes))
    slides_text: list[str] = []

    for slide_num, slide in enumerate(prs.slides, start=1):
        lines: list[str] = [f"--- Slide {slide_num} ---"]
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = para.text.strip()
                    if line:
                        lines.append(line)
        slides_text.append("\n".join(lines))

    text = "\n\n".join(slides_text).strip()

    if len(text) < 100:
        raise ValueError(
            f"PPTX extraction yielded only {len(text)} characters — "
            "presentation may be empty or contain only images."
        )

    return text
